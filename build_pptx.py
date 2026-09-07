"""
Build ocm_presentation.pptx — 15-slide PowerPoint deck.

Mirrors the 9-chapter walkthrough notebook narrative:
  1  Title
  2  The Problem — Chapter 1
  3  Setup — Label Shift & PCA — Chapter 2
  4  Setup — Element Usage — Chapter 2
  5  How We Measure Success — Chapter 3
  6  Baseline + Why Naive Merging Fails — Chapter 4
  7  DRST — Filtering by Chemical Similarity — Chapter 5
  8  KMM — Soft Weights Instead of a Hard Filter — Chapter 6
  9  Prior Feature Transfer Pipeline — Chapter 7
 10  Results — All Five Methods — Chapter 8
 11  SHAP — Feature Importance Beeswarm — Chapter 8
 12  SHAP — Ceiling Effect & Model Limits — Chapter 8
 13  Statistical Rigour & Honesty (10-seed, QN trade-off) — Chapter 8
 14  Critical Analysis — Per-Method Review — Chapter 9
 15  What's Next — Priority-Ordered — Chapter 9
 16  Summary

Run: python build_pptx.py
Produces: ocm_presentation.pptx  (and ocm_presentation.pdf via LibreOffice)
"""

import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = pathlib.Path(__file__).parent

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x1E, 0x42)
ACCENT = RGBColor(0x00, 0x78, 0xC8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF6, 0xF8)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
RED    = RGBColor(0xDC, 0x3C, 0x3C)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
DGRAY  = RGBColor(0x50, 0x59, 0x6E)
LGRAY  = RGBColor(0xD0, 0xD5, 0xDF)
BLACK  = RGBColor(0x1A, 0x1A, 0x2E)
GDARK  = RGBColor(0x1B, 0x5E, 0x20)
RDARK  = RGBColor(0x8B, 0x1A, 0x1A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line; shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_img(slide, fname, x, y, w, h=None):
    path = ROOT / fname
    if not path.exists():
        return
    if h:
        slide.shapes.add_picture(str(path), x, y, w, h)
    else:
        slide.shapes.add_picture(str(path), x, y, w)

def nav_bar(slide, current, total):
    add_rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=NAVY)
    add_text(slide, f"OCM Dataset Integration   |   Slide {current} / {total}",
             Inches(0.2), H - Inches(0.36), W - Inches(0.4), Inches(0.34),
             size=9, color=RGBColor(0xCC, 0xCC, 0xCC))

def slide_title_bar(slide, title, subtitle=None, chapter=None):
    has_sub = bool(subtitle) or bool(chapter)
    bar_h = Inches(1.2) if has_sub else Inches(0.95)
    add_rect(slide, 0, 0, W, bar_h, fill=NAVY)
    add_text(slide, title, Inches(0.4), Inches(0.18), W - Inches(0.8), Inches(0.55),
             size=22, bold=True, color=WHITE)
    if has_sub:
        sub_text = subtitle or ""
        if chapter:
            sub_text = (chapter + "   |   " + sub_text) if sub_text else chapter
        add_text(slide, sub_text, Inches(0.4), Inches(0.72), W - Inches(0.8), Inches(0.4),
                 size=11, color=RGBColor(0xBB, 0xCC, 0xEE), italic=True)

def bullets(slide, items, x, y, w, h, size=12, spacing=0.34, color=BLACK):
    for i, item in enumerate(items):
        add_text(slide, "•  " + item, x, y + Inches(i * spacing), w, Inches(spacing + 0.05),
                 size=size, color=color)

def badge(slide, text, x, y, w, h, fill, text_color=WHITE, size=11):
    add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, text, x + Inches(0.05), y + Inches(0.02),
             w - Inches(0.1), h - Inches(0.04),
             size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def callout(slide, text, x, y, w, h, fill=LIGHT, border=ACCENT,
            text_color=DGRAY, size=11, italic=True):
    add_rect(slide, x, y, w, h, fill=fill, line=border)
    add_text(slide, text, x + Inches(0.12), y + Inches(0.08),
             w - Inches(0.24), h - Inches(0.16),
             size=size, color=text_color, italic=italic)

TOTAL = 16

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=NAVY)
add_rect(s, 0, Inches(2.95), W, Inches(0.06), fill=ACCENT)

add_text(s, "JAIST — Taniike Lab",
         Inches(1.2), Inches(1.4), Inches(10.9), Inches(0.5),
         size=13, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
add_text(s, "Integrating OCM Literature Data\ninto Lab-Scale ML Models",
         Inches(1.2), Inches(1.9), Inches(10.9), Inches(1.6),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "A Walkthrough of Five Transfer-Learning Approaches",
         Inches(1.2), Inches(3.55), Inches(10.9), Inches(0.6),
         size=16, color=RGBColor(0x99, 0xBB, 0xEE), align=PP_ALIGN.CENTER, italic=True)
add_rect(s, Inches(4.5), Inches(4.35), Inches(4.33), Inches(0.04),
         fill=RGBColor(0x44, 0x66, 0xAA))
add_text(s, "Partha Nupam Nagar",
         Inches(1.2), Inches(4.55), Inches(10.9), Inches(0.45),
         size=15, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, "June 2026",
         Inches(1.2), Inches(5.0), Inches(10.9), Inches(0.4),
         size=13, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
nav_bar(s, 1, TOTAL)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem (two datasets + three differences)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "The Problem — Two Datasets, Three Systematic Differences", chapter="Chapter 1")
nav_bar(s, 2, TOTAL)

# Left dataset card
add_rect(s, Inches(0.4), Inches(1.2), Inches(6.0), Inches(2.6),
         fill=RGBColor(0xE8, 0xF4, 0xFF), line=ACCENT)
add_rect(s, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.42), fill=ACCENT)
add_text(s, "Our Lab Data", Inches(0.55), Inches(1.23),
         Inches(5.7), Inches(0.36), size=13, bold=True, color=WHITE)
rows_l = [("Samples", "89,074"), ("Year", "2025"),
          ("Preparation", "Impregnation only"),
          ("Mean Y(C2)", "5.25%"), ("Conditions", "Fixed, one lab")]
for i, (k, v) in enumerate(rows_l):
    add_text(s, k, Inches(0.55), Inches(1.7) + Inches(i * 0.38),
             Inches(2.4), Inches(0.34), size=11, bold=True, color=DGRAY)
    add_text(s, v, Inches(2.95), Inches(1.7) + Inches(i * 0.38),
             Inches(3.3), Inches(0.34), size=11, color=BLACK)

# Right dataset card
add_rect(s, Inches(6.95), Inches(1.2), Inches(6.0), Inches(2.6),
         fill=RGBColor(0xFF, 0xF5, 0xE8), line=AMBER)
add_rect(s, Inches(6.95), Inches(1.2), Inches(6.0), Inches(0.42), fill=AMBER)
add_text(s, "Published Literature Data (≤2019)", Inches(7.1), Inches(1.23),
         Inches(5.7), Inches(0.36), size=13, bold=True, color=WHITE)
rows_r = [("Samples", "3,852"), ("Year", "1982 – 2019"),
          ("Preparation", "15+ methods"),
          ("Mean Y(C2)", "8.67%"), ("Conditions", "40 years of diverse labs")]
for i, (k, v) in enumerate(rows_r):
    add_text(s, k, Inches(7.1), Inches(1.7) + Inches(i * 0.38),
             Inches(2.4), Inches(0.34), size=11, bold=True, color=DGRAY)
    add_text(s, v, Inches(9.55), Inches(1.7) + Inches(i * 0.38),
             Inches(3.3), Inches(0.34), size=11, color=BLACK)

# Central question
add_text(s,
         "Can we use literature to improve our model — without hurting it?",
         Inches(0.4), Inches(3.95), Inches(12.5), Inches(0.4),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, italic=True)

# Three differences
diffs = [
    ("1. Label shift",
     "Literature mean is 3.42 pp higher (t = −32.2, p < 10⁻²⁰⁰) — a systematic offset from publication bias and optimised conditions.",
     RED),
    ("2. Covariate shift",
     "78.5% of literature samples describe chemistry our lab has never tested.",
     AMBER),
    ("3. Publication bias",
     "Literature is skewed, not just shifted — failed experiments are rarely reported, over-weighting high yields.",
     ACCENT),
]
for i, (head, body, color) in enumerate(diffs):
    y0 = Inches(4.5) + Inches(i * 0.62)
    badge(s, head, Inches(0.4), y0, Inches(2.7), Inches(0.42), color, size=12)
    add_text(s, body, Inches(3.25), y0 + Inches(0.04),
             Inches(9.7), Inches(0.4), size=11, color=DGRAY)

# Why ML? callout
callout(s,
        "Why ML? The catalyst search space (65 elements × varying loadings × temperature) is too large to screen physically. "
        "A surrogate model guides experiments toward high-yield compositions.",
        Inches(0.4), Inches(6.42), Inches(12.5), Inches(0.68),
        fill=LIGHT, border=NAVY, text_color=DGRAY, size=10, italic=False)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Setup: Label Shift & PCA (Chapter 2a)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Setup — Label Shift & Covariate Shift", chapter="Chapter 2")
nav_bar(s, 3, TOTAL)

add_img(s, "fig_walkthrough_label_pca.png", Inches(0.3), Inches(1.1), Inches(9.0))

# Right commentary
add_rect(s, Inches(9.6), Inches(1.1), Inches(3.4), Inches(0.4), fill=NAVY)
add_text(s, "Two panels, two shifts", Inches(9.75), Inches(1.13),
         Inches(3.1), Inches(0.34), size=12, bold=True, color=WHITE)
bullets(s, [
    "Left: yield KDE — blue peaks 3–4%",
    "Orange peaks 8–9% (3.42 pp gap)",
    "Red band = systematic gap (t = −32.2)",
    "Orange right tail = publication-bias skew",
    "Right: PCA of 67 features → 2D",
    "Orange outside blue = 78.5% OOD",
], Inches(9.6), Inches(1.62), Inches(3.4), Inches(0.4), size=11, spacing=0.40, color=BLACK)

callout(s,
        "Both signals — yield level AND chemistry composition — warn against naive concatenation (confirmed slide 6).",
        Inches(9.6), Inches(4.15), Inches(3.4), Inches(0.9),
        fill=RGBColor(0xFF, 0xF0, 0xD8), border=AMBER,
        text_color=RGBColor(0x6B, 0x3A, 0x00), size=10, italic=False)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Setup: Element Usage (Chapter 2b)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Setup — Which Elements Differ?", chapter="Chapter 2")
nav_bar(s, 4, TOTAL)

add_img(s, "fig_element_usage.png", Inches(0.3), Inches(1.1), Inches(9.0))

# Right commentary
add_rect(s, Inches(9.6), Inches(1.1), Inches(3.4), Inches(0.4), fill=ACCENT)
add_text(s, "Reading the chart", Inches(9.75), Inches(1.13),
         Inches(3.1), Inches(0.34), size=12, bold=True, color=WHITE)
bullets(s, [
    "Top 15 elements by usage frequency",
    "Different dominant elements each side",
    "Different active phases & promoters",
    "Our lab: focused 2025 programme",
    "Literature: 40 yrs, many groups",
], Inches(9.6), Inches(1.62), Inches(3.4), Inches(0.4), size=11, spacing=0.40, color=BLACK)

callout(s,
        "This is what '78.5% out-of-distribution' looks like in plain terms — literally different elements and preparation methods.",
        Inches(9.6), Inches(3.85), Inches(3.4), Inches(0.9),
        fill=RGBColor(0xE8, 0xF5, 0xFF), border=ACCENT,
        text_color=DGRAY, size=10, italic=False)

add_text(s,
         "Selective filtering methods keep only literature whose element profile resembles our experimental programme.",
         Inches(9.6), Inches(4.9), Inches(3.4), Inches(0.65),
         size=10, color=DGRAY, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — How We Measure Success (Chapter 3)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "How We Measure Success — Catalyst-Grouped 5-Fold CV", chapter="Chapter 3")
nav_bar(s, 5, TOTAL)

# Left: rule + reasoning
add_rect(s, Inches(0.4), Inches(1.2), Inches(7.0), Inches(0.42), fill=NAVY)
add_text(s, "The validation rule", Inches(0.55), Inches(1.23),
         Inches(6.7), Inches(0.36), size=13, bold=True, color=WHITE)
bullets(s, [
    "The 89,074 rows are only 917 DISTINCT CATALYSTS",
    "Folds split by catalyst — all rows of a catalyst in one fold",
    "Validation fold = catalysts the model has never seen",
    "Literature never appears in validation",
], Inches(0.45), Inches(1.7), Inches(6.9), Inches(0.4), size=12, spacing=0.4)

add_rect(s, Inches(0.4), Inches(3.65), Inches(7.0), Inches(0.42), fill=ACCENT)
add_text(s, "Why grouped, and not by row", Inches(0.55), Inches(3.68),
         Inches(6.7), Inches(0.36), size=13, bold=True, color=WHITE)
bullets(s, [
    "A row split puts the SAME catalyst in train and test",
    "The model then recalls catalysts instead of predicting new ones",
    "Prof. Taniike raised exactly this — grouped CV is now the default",
], Inches(0.45), Inches(4.15), Inches(6.9), Inches(0.4), size=12, spacing=0.4)

# Right: metric box + fold diagram
add_rect(s, Inches(7.7), Inches(1.2), Inches(5.3), Inches(1.4),
         fill=RGBColor(0xE8, 0xF5, 0xE9), line=GREEN)
add_text(s, "Primary metric: catalyst-level",
         Inches(7.85), Inches(1.3), Inches(5.0), Inches(0.4),
         size=14, bold=True, color=GDARK)
add_text(s, "Spearman on each catalyst's best yield, and enrichment@10%.\nRow RMSE is secondary — 19.9% of its variance is unreachable.",
         Inches(7.85), Inches(1.75), Inches(5.0), Inches(0.75),
         size=11, color=DGRAY)

# Fold visualization
add_text(s, "5-fold split by CATALYST (917 groups)",
         Inches(7.7), Inches(2.85), Inches(5.3), Inches(0.4),
         size=12, bold=True, color=NAVY)

fold_colors = [ACCENT, ACCENT, RED, ACCENT, ACCENT]
fold_labels = ["~183\ncatalysts", "~183\ncatalysts", "~183\nUNSEEN",
               "~183\ncatalysts", "~183\ncatalysts"]
for i, (col, lab) in enumerate(zip(fold_colors, fold_labels)):
    x = Inches(7.7) + Inches(i * 1.05)
    add_rect(s, x, Inches(3.35), Inches(1.0), Inches(0.85), fill=col)
    add_text(s, lab, x, Inches(3.4), Inches(1.0), Inches(0.8),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "+ literature (training only, never validation)",
         Inches(7.7), Inches(4.4), Inches(5.3), Inches(0.4),
         size=11, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

callout(s,
        "Rotate the held-out fold 5 times. Every number on the following slides uses this protocol; "
        "row-level figures appear only where they are labelled as historical.",
        Inches(7.7), Inches(5.0), Inches(5.3), Inches(0.8),
        fill=LIGHT, border=NAVY, text_color=DGRAY, size=10, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Baseline + Naive Merge (Chapter 4)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Baseline + Why Naive Merging Fails", chapter="Chapter 4")
nav_bar(s, 6, TOTAL)

# Left: Step 1 + Step 2 tables
add_text(s, "Step 1 — establish baseline  (ROW-LEVEL CV, historical)",
         Inches(0.4), Inches(1.3), Inches(6.2), Inches(0.4),
         size=13, bold=True, color=NAVY)
add_rect(s, Inches(0.4), Inches(1.8), Inches(6.2), Inches(0.42), fill=NAVY)
add_text(s, "Method", Inches(0.55), Inches(1.83), Inches(4.0), Inches(0.36),
         size=12, bold=True, color=WHITE)
add_text(s, "CV RMSE", Inches(4.7), Inches(1.83), Inches(1.8), Inches(0.36),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(0.4), Inches(2.22), Inches(6.2), Inches(0.5),
         fill=LIGHT, line=LGRAY)
add_text(s, "Baseline (lab data only)",
         Inches(0.55), Inches(2.27), Inches(4.0), Inches(0.42),
         size=12, color=BLACK)
add_text(s, "2.133", Inches(4.7), Inches(2.27), Inches(1.8), Inches(0.42),
         size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, "Row-level protocol — kept only for comparison with the published numbers. "
            "Under catalyst-grouped CV the baseline is 2.9425.",
         Inches(0.4), Inches(2.85), Inches(6.2), Inches(0.5),
         size=10, italic=True, color=DGRAY)

add_text(s, "Step 2 — try the obvious",
         Inches(0.4), Inches(3.55), Inches(6.2), Inches(0.4),
         size=14, bold=True, color=NAVY)
add_rect(s, Inches(0.4), Inches(4.05), Inches(6.2), Inches(0.42), fill=NAVY)
add_text(s, "Method", Inches(0.55), Inches(4.08), Inches(4.0), Inches(0.36),
         size=12, bold=True, color=WHITE)
add_text(s, "CV RMSE", Inches(4.7), Inches(4.08), Inches(1.8), Inches(0.36),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(0.4), Inches(4.47), Inches(6.2), Inches(0.5),
         fill=LIGHT, line=LGRAY)
add_text(s, "Baseline", Inches(0.55), Inches(4.52), Inches(4.0), Inches(0.42),
         size=12, color=BLACK)
add_text(s, "2.133", Inches(4.7), Inches(4.52), Inches(1.8), Inches(0.42),
         size=12, color=BLACK, align=PP_ALIGN.CENTER)
add_rect(s, Inches(0.4), Inches(4.97), Inches(6.2), Inches(0.5),
         fill=RGBColor(0xFF, 0xEB, 0xEB), line=RED)
add_text(s, "Naive merge (3,852 lit rows)",
         Inches(0.55), Inches(5.02), Inches(4.0), Inches(0.42),
         size=12, bold=True, color=BLACK)
add_text(s, "2.241", Inches(4.7), Inches(5.02), Inches(1.8), Inches(0.42),
         size=12, bold=True, color=RDARK, align=PP_ALIGN.CENTER)
add_text(s, "+5.1% worse than baseline (row-level). Re-tested at catalyst level, a direct "
            "merge again sits below the composition-only control (0.7577 vs 0.7606).",
         Inches(0.4), Inches(5.6), Inches(6.2), Inches(0.6),
         size=11, bold=True, italic=True, color=RDARK)

# Right: explanation
add_rect(s, Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.42), fill=ACCENT)
add_text(s, "Why did it fail?", Inches(7.15), Inches(1.33),
         Inches(5.7), Inches(0.36), size=13, bold=True, color=WHITE)
add_rect(s, Inches(7.0), Inches(1.72), Inches(6.0), Inches(2.7),
         fill=LIGHT, line=LGRAY)
add_text(s,
         "Model sees two catalysts with nearly identical composition:",
         Inches(7.15), Inches(1.85), Inches(5.7), Inches(0.35),
         size=12, color=BLACK)
bullets(s, [
    "One labelled 5.2%  (our lab)",
    "One labelled 8.7%  (literature)",
], Inches(7.25), Inches(2.25), Inches(5.6), Inches(0.4), size=12, spacing=0.4)
add_text(s,
         "No feature explains the gap — flow rate, gas ratio, etc. are absent. The model hedges and predicts both badly.",
         Inches(7.15), Inches(3.1), Inches(5.7), Inches(1.2),
         size=12, color=DGRAY)

callout(s,
        "The 3.42 pp offset is a SYSTEMATIC signal the model cannot account for — and incorrectly tries to fit. "
        "Mixing the two datasets can perform WORSE than using ours alone.\n\n"
        "⇒ We need to be much more selective.",
        Inches(7.0), Inches(4.35), Inches(6.0), Inches(1.9),
        fill=RGBColor(0xFF, 0xF0, 0xD8), border=AMBER,
        text_color=RGBColor(0x6B, 0x3A, 0x00), size=11, italic=False)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DRST (Chapter 5)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "DRST — Filtering by Chemical Similarity", chapter="Chapter 5")
nav_bar(s, 7, TOTAL)

add_img(s, "fig_drst_scores.png", Inches(0.3), Inches(1.25), Inches(8.0))
add_text(s, "Histogram of P(ours | x) for every literature sample.  Dashed lines = candidate thresholds.",
         Inches(0.3), Inches(4.8), Inches(8.0), Inches(0.4),
         size=10, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

# Right: how it works + sweep table
add_rect(s, Inches(8.5), Inches(1.25), Inches(4.5), Inches(0.42), fill=NAVY)
add_text(s, "How it works", Inches(8.65), Inches(1.28),
         Inches(4.2), Inches(0.36), size=13, bold=True, color=WHITE)
bullets(s, [
    "Train logistic classifier: ours=1, lit=0",
    "Score each lit sample: P(ours | x)",
    "Keep samples scoring ≥ τ",
], Inches(8.55), Inches(1.75), Inches(4.3), Inches(0.4), size=11, spacing=0.4)
callout(s,
        "P(ours | x) is a similarity score: high = looks like our chemistry, "
        "low = foreign. We keep only literature scoring above τ.",
        Inches(8.5), Inches(3.05), Inches(4.5), Inches(0.65),
        fill=RGBColor(0xE8, 0xF5, 0xFF), border=ACCENT, text_color=DGRAY, size=9, italic=False)

# Sweep table
add_text(s, "Threshold sweep",
         Inches(8.5), Inches(3.8), Inches(4.5), Inches(0.4),
         size=13, bold=True, color=NAVY)
sweep_rows = [
    ("τ = 0.10", "1168 kept", "2.068", False),
    ("τ = 0.20", "987 kept",  "2.034", False),
    ("best τ", "~200 kept",  "2.127", False),
    ("τ = 0.40", "543 kept",  "2.045", False),
]
add_rect(s, Inches(8.5), Inches(4.25), Inches(4.5), Inches(0.4), fill=NAVY)
for j, (h, x, w) in enumerate(zip(["τ", "kept", "RMSE"],
                                    [Inches(8.65), Inches(9.95), Inches(11.6)],
                                    [Inches(1.25), Inches(1.6), Inches(0.9)])):
    add_text(s, h, x, Inches(4.28), w, Inches(0.34),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, (tau, kept, rmse, best) in enumerate(sweep_rows):
    y0 = Inches(4.65) + Inches(i * 0.4)
    bg = RGBColor(0xE8, 0xF5, 0xE9) if best else LIGHT
    add_rect(s, Inches(8.5), y0, Inches(4.5), Inches(0.4), fill=bg, line=LGRAY)
    add_text(s, tau, Inches(8.65), y0 + Inches(0.04), Inches(1.25), Inches(0.34),
             size=11, bold=best, color=GDARK if best else BLACK, align=PP_ALIGN.CENTER)
    add_text(s, kept, Inches(9.95), y0 + Inches(0.04), Inches(1.6), Inches(0.34),
             size=11, bold=best, color=GDARK if best else BLACK, align=PP_ALIGN.CENTER)
    add_text(s, rmse, Inches(11.6), y0 + Inches(0.04), Inches(0.9), Inches(0.34),
             size=11, bold=best, color=GDARK if best else BLACK, align=PP_ALIGN.CENTER)

callout(s,
        "Unsolved: hard cutoff is arbitrary — a sample at 0.29 is fully discarded; 0.31 gets full weight.",
        Inches(8.5), Inches(6.35), Inches(4.5), Inches(0.72),
        fill=RGBColor(0xFF, 0xF0, 0xD8), border=AMBER,
        text_color=RGBColor(0x6B, 0x3A, 0x00), size=10)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — KMM (Chapter 6)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "KMM — Soft Weights Instead of a Hard Filter", chapter="Chapter 6")
nav_bar(s, 8, TOTAL)

add_img(s, "fig_kmm_weights.png", Inches(0.3), Inches(1.25), Inches(8.65))

# Right
add_rect(s, Inches(9.2), Inches(1.25), Inches(3.8), Inches(0.42), fill=NAVY)
add_text(s, "How it works", Inches(9.35), Inches(1.28),
         Inches(3.5), Inches(0.36), size=13, bold=True, color=WHITE)
bullets(s, [
    "Each lit sample gets weight wᵢ ∈ [0, 10]",
    "RBF kernel: K(x,x′) = exp(−‖x−x′‖²/2σ²)",
    "Solve QP:  ½ wᵀK_ss w − κᵀw",
    "σ via median heuristic (parameter-free)",
], Inches(9.25), Inches(1.75), Inches(3.7), Inches(0.4), size=11, spacing=0.38)
callout(s,
        "KMM picks weights so the weighted literature cloud matches lab data cloud. "
        "Samples inside our chemistry earn weight; those outside fade to zero.",
        Inches(9.2), Inches(3.35), Inches(3.8), Inches(0.62),
        fill=RGBColor(0xE8, 0xF5, 0xFF), border=ACCENT, text_color=DGRAY, size=9, italic=False)

# Result table
add_text(s, "Result",
         Inches(9.2), Inches(4.08), Inches(3.8), Inches(0.4),
         size=13, bold=True, color=NAVY)
results_rows = [
    ("KMM CV RMSE",        "2.261"),
    ("Δ vs baseline",      "+6.0% (worse)"),
    ("Near-zero weights",  "78.5%"),
    ("Agreement w/ DRST",  "r = 0.79"),
]
for i, (k, v) in enumerate(results_rows):
    y0 = Inches(4.5) + Inches(i * 0.4)
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(9.2), y0, Inches(3.8), Inches(0.4), fill=bg, line=LGRAY)
    add_text(s, k, Inches(9.35), y0 + Inches(0.05), Inches(2.4), Inches(0.34),
             size=11, color=DGRAY)
    add_text(s, v, Inches(11.65), y0 + Inches(0.05), Inches(1.3), Inches(0.34),
             size=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

callout(s,
        "Reassuring: two unrelated methods flag the same 78.5% of samples ⇒ real signal, not artifact.\n"
        "Still unsolved: literature labels remain in training, so the 3.42 pp offset still bleeds in.",
        Inches(9.2), Inches(6.2), Inches(3.8), Inches(0.85),
        fill=LIGHT, border=ACCENT, text_color=DGRAY, size=10)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Prior Feature Transfer Pipeline (Chapter 7)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Prior Feature Transfer — The Method That Exposed the Leak", chapter="Chapter 7")
nav_bar(s, 9, TOTAL)

# Stage labels
add_text(s, "STAGE 1 — Pre-train on Literature Chemistry",
         Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.38),
         size=13, bold=True, color=ACCENT)
add_text(s, "STAGE 2 — Fine-tune on Our Labels Only",
         Inches(0.3), Inches(3.9), Inches(12.7), Inches(0.38),
         size=13, bold=True, color=GREEN)

def pipeline_box(slide, x, y, w, h, title, subtitle, fill_c, text_c=WHITE):
    add_rect(slide, x, y, w, h, fill=fill_c, line=RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, title, x + Inches(0.1), y + Inches(0.08),
             w - Inches(0.2), Inches(0.4), size=12, bold=True,
             color=text_c, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, x + Inches(0.1), y + Inches(0.5),
                 w - Inches(0.2), Inches(0.5),
                 size=10, color=text_c, align=PP_ALIGN.CENTER)

def arrow_h(slide, x, y, length, color=ACCENT):
    line = slide.shapes.add_connector(1, x, y, x + length, y)
    line.line.color.rgb = color; line.line.width = Pt(2.5)

bw = Inches(2.5); bh = Inches(1.0)
by1 = Inches(1.7); by2 = Inches(4.4)
gap = Inches(0.5)

# Stage 1 row
pipeline_box(s, Inches(0.4), by1, bw, bh,
             "Literature", "3,852 samples\n(DRST filtered)",
             RGBColor(0xCC, 0xE0, 0xFF), text_c=NAVY)
arrow_h(s, Inches(0.4) + bw, by1 + bh/2, gap)
pipeline_box(s, Inches(0.4) + bw + gap, by1, bw, bh,
             "Stage 1", "XGBoost\npre-train",
             ACCENT, text_c=WHITE)
arrow_h(s, Inches(0.4) + 2*bw + gap, by1 + bh/2, gap)
pipeline_box(s, Inches(0.4) + 2*(bw + gap), by1, bw, bh,
             "lit_prior", "prediction\n(68th feature)",
             AMBER, text_c=WHITE)

# Stage 2 row
pipeline_box(s, Inches(0.4), by2, bw, bh,
             "Our Lab Data", "89,074 samples",
             RGBColor(0xCC, 0xE8, 0xD0), text_c=GDARK)
arrow_h(s, Inches(0.4) + bw, by2 + bh/2, gap, color=GREEN)
pipeline_box(s, Inches(0.4) + bw + gap, by2, bw, bh,
             "Append", "68th feature",
             RGBColor(0xE0, 0xE8, 0xF8), text_c=NAVY)
arrow_h(s, Inches(0.4) + 2*bw + gap, by2 + bh/2, gap, color=GREEN)
pipeline_box(s, Inches(0.4) + 2*(bw + gap), by2, bw, bh,
             "Stage 2", "LightGBM\nour labels only",
             GREEN, text_c=WHITE)
arrow_h(s, Inches(0.4) + 3*bw + 2*gap, by2 + bh/2, gap, color=GREEN)
pipeline_box(s, Inches(0.4) + 3*(bw + gap), by2, bw, bh,
             "Prediction", "Y(C₂) %",
             NAVY, text_c=WHITE)

# Vertical arrow from prior to append
line = s.shapes.add_connector(1,
                              Inches(0.4) + 2*(bw + gap) + bw/2, by1 + bh,
                              Inches(0.4) + bw + gap + bw/2, by2)
line.line.color.rgb = AMBER; line.line.width = Pt(2.5)

# Key property callout
add_rect(s, Inches(0.4), Inches(5.65), Inches(12.5), Inches(1.45),
         fill=LIGHT, line=NAVY)
add_text(s, "The design was sound. The evaluation was not.",
         Inches(0.55), Inches(5.72), Inches(12.2), Inches(0.4),
         size=13, bold=True, color=NAVY)
add_text(s,
         "Literature enters as a FEATURE, never as a label, so the 3.42 pp offset cannot corrupt the training loss. That reasoning still holds. "
         "But Stage 1 was trained on literature TOGETHER WITH the lab training rows, and under a row-level split those rows included the test "
         "catalysts — so the prior feature partly carried each test catalyst's own measured yields. The gain was recall, not prediction. "
         "Ruling this out is what produced the protocol on slide 5; the mechanism is named in ocm_eval.stage1_data().",
         Inches(0.55), Inches(6.12), Inches(12.2), Inches(0.9),
         size=10, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Results (Chapter 8)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Results — The Same Models Under Two Protocols", chapter="Chapter 8")
nav_bar(s, 10, TOTAL)

add_img(s, "fig_protocol_comparison.png", Inches(0.3), Inches(1.25), Inches(8.3))

# Right: CV table + OOD table
add_text(s, "Row-level CV — HISTORICAL, leaked",
         Inches(8.8), Inches(1.3), Inches(4.2), Inches(0.4),
         size=14, bold=True, color=RDARK)
cv_rows = [
    ("Baseline",        "2.1184", "—",      False, False),
    ("PFT (filtered)",  "1.9120", "−9.7%",  True,  False),
    ("PFT (all lit.)",  "1.9194", "−9.4%",  False, False),
]
add_rect(s, Inches(8.8), Inches(1.75), Inches(4.2), Inches(0.4), fill=NAVY)
for j, (h, x, w) in enumerate(zip(["Method", "RMSE", "Δ"],
                                    [Inches(8.9), Inches(10.8), Inches(11.9)],
                                    [Inches(1.85), Inches(1.0), Inches(1.1)])):
    add_text(s, h, x, Inches(1.78), w, Inches(0.34),
             size=11, bold=True, color=WHITE)
for i, (m, r, d, best, worst) in enumerate(cv_rows):
    y0 = Inches(2.15) + Inches(i * 0.4)
    bg = RGBColor(0xE8, 0xF5, 0xE9) if best else \
         RGBColor(0xFF, 0xEB, 0xEB) if worst else (LIGHT if i % 2 == 0 else WHITE)
    add_rect(s, Inches(8.8), y0, Inches(4.2), Inches(0.4), fill=bg, line=LGRAY)
    fc = GDARK if best else (RDARK if worst else BLACK)
    add_text(s, m, Inches(8.9), y0 + Inches(0.04), Inches(1.85), Inches(0.34),
             size=11, bold=best, color=fc)
    add_text(s, r, Inches(10.8), y0 + Inches(0.04), Inches(1.0), Inches(0.34),
             size=11, bold=best, color=fc, align=PP_ALIGN.CENTER)
    add_text(s, d, Inches(11.9), y0 + Inches(0.04), Inches(1.1), Inches(0.34),
             size=11, bold=best, color=fc, align=PP_ALIGN.CENTER)

# Catalyst-grouped — the default protocol, and the inversion
add_text(s, "Catalyst-grouped CV — DEFAULT",
         Inches(8.8), Inches(3.95), Inches(4.2), Inches(0.4),
         size=14, bold=True, color=GDARK)
add_rect(s, Inches(8.8), Inches(4.4), Inches(4.2), Inches(0.4), fill=NAVY)
for h, x, w in zip(["Method", "RMSE", "Δ"],
                   [Inches(8.9), Inches(10.8), Inches(11.9)],
                   [Inches(1.85), Inches(1.0), Inches(1.1)]):
    add_text(s, h, x, Inches(4.43), w, Inches(0.34), size=11, bold=True, color=WHITE)
grouped_rows = [("Baseline", "2.9425", "—", False),
                ("PFT (filtered)", "2.9955", "+1.8%", True),
                ("PFT (all lit.)", "2.9817", "+1.3%", True)]
for i, (m, r, dlt, worse) in enumerate(grouped_rows):
    y0 = Inches(4.8) + Inches(i * 0.4)
    bg = RGBColor(0xFF, 0xEB, 0xEB) if worse else LIGHT
    add_rect(s, Inches(8.8), y0, Inches(4.2), Inches(0.4), fill=bg, line=LGRAY)
    fc = RDARK if worse else BLACK
    add_text(s, m, Inches(8.9), y0 + Inches(0.04), Inches(1.85), Inches(0.34),
             size=11, color=fc)
    add_text(s, r, Inches(10.8), y0 + Inches(0.04), Inches(1.0), Inches(0.34),
             size=11, color=fc, align=PP_ALIGN.CENTER)
    add_text(s, dlt, Inches(11.9), y0 + Inches(0.04), Inches(1.1), Inches(0.34),
             size=11, color=fc, align=PP_ALIGN.CENTER)

callout(s,
        "Same models, same data, two protocols. The published −9.7% becomes +1.8% WORSE once a "
        "catalyst cannot appear in both train and test. The difference is leakage, not modelling.",
        Inches(8.8), Inches(6.1), Inches(4.2), Inches(0.95),
        fill=RGBColor(0xFF, 0xF4, 0xE5), border=RGBColor(0xE6, 0x7E, 0x22),
        text_color=DGRAY, size=10)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SHAP Beeswarm standalone (Chapter 8)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "SHAP — Feature Importance Beeswarm", chapter="Chapter 8")
nav_bar(s, 11, TOTAL)

add_img(s, "fig_shap_beeswarm.png", Inches(0.3), Inches(1.2), Inches(9.5))
add_text(s,
         "TreeExplainer on 3,000-sample subsample.  Each row = feature ranked by mean |SHAP|.  "
         "Colour = feature value (red=high);  x-position = prediction impact.",
         Inches(0.3), Inches(5.5), Inches(9.5), Inches(0.5),
         size=9, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

# Right: top findings
add_rect(s, Inches(10.1), Inches(1.2), Inches(2.9), Inches(0.42), fill=NAVY)
add_text(s, "Key signals", Inches(10.25), Inches(1.23),
         Inches(2.65), Inches(0.36), size=12, bold=True, color=WHITE)
bullets(s, [
    "#1: lit_prior — transfer used",
    "Temp ↑ → yield ↑ (correct OCM)",
    "Ba, Mn, La, Ce: net positive",
    "Li, K: mixed / suppressive",
], Inches(10.1), Inches(1.74), Inches(2.85), Inches(0.4), size=11, spacing=0.44, color=BLACK)

callout(s,
        "The model uses real OCM chemistry — transfer learning is not dead weight.",
        Inches(10.1), Inches(3.6), Inches(2.9), Inches(0.75),
        fill=RGBColor(0xE8, 0xF5, 0xFF), border=ACCENT, text_color=DGRAY, size=10)

add_text(s, "See next slide for residual breakdown →",
         Inches(10.1), Inches(4.6), Inches(2.9), Inches(0.5),
         size=10, italic=True, color=DGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SHAP Ceiling Effect & Findings (Chapter 8)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "SHAP — Ceiling Effect & Model Limits", chapter="Chapter 8")
nav_bar(s, 12, TOTAL)

# Left: residual table
add_text(s, "Residual breakdown by yield range",
         Inches(0.4), Inches(1.3), Inches(5.8), Inches(0.4),
         size=14, bold=True, color=NAVY)
add_rect(s, Inches(0.4), Inches(1.75), Inches(5.8), Inches(0.42), fill=NAVY)
add_text(s, "Y(C₂) range", Inches(0.55), Inches(1.78), Inches(3.1), Inches(0.34),
         size=11, bold=True, color=WHITE)
add_text(s, "RMSE", Inches(3.8), Inches(1.78), Inches(2.2), Inches(0.34),
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
res_rows = [
    ("0 – 6%",   "1.43 – 1.48", False),
    ("6 – 10%",  "1.95",        False),
    ("10 – 15%", "2.58",        False),
    ("> 15%",    "4.70  (−4.2 pp bias)", True),
]
for i, (rng, rmse, bad) in enumerate(res_rows):
    y0 = Inches(2.17) + Inches(i * 0.52)
    bg = RGBColor(0xFF, 0xEB, 0xEB) if bad else (LIGHT if i % 2 == 0 else WHITE)
    add_rect(s, Inches(0.4), y0, Inches(5.8), Inches(0.52), fill=bg, line=LGRAY)
    fc = RDARK if bad else BLACK
    add_text(s, rng, Inches(0.55), y0 + Inches(0.08), Inches(3.1), Inches(0.38),
             size=12, bold=bad, color=fc)
    add_text(s, rmse, Inches(3.8), y0 + Inches(0.08), Inches(2.2), Inches(0.38),
             size=12, bold=bad, color=fc, align=PP_ALIGN.CENTER)

callout(s,
        "Root cause: GHSV, CH₄/O₂ ratio, pressure absent from the feature set — "
        "a data limit, not a model failure.",
        Inches(0.4), Inches(4.4), Inches(5.8), Inches(0.85),
        fill=RGBColor(0xFF, 0xF0, 0xD8), border=AMBER,
        text_color=RGBColor(0x6B, 0x3A, 0x00), size=10, italic=False)

# Right: top features + ceiling insight
add_rect(s, Inches(6.7), Inches(1.3), Inches(6.3), Inches(0.42), fill=ACCENT)
add_text(s, "What SHAP confirmed",
         Inches(6.85), Inches(1.33), Inches(6.0), Inches(0.36),
         size=13, bold=True, color=WHITE)
bullets(s, [
    "#1: lit_prior_prediction — transfer learning actively used",
    "Temperature: high T → yield up (correct OCM physics)",
    "Ba, Mn, La, Ce — known active phases and promoters",
    "Li, K — mixed / suppressive (can over-reduce surface)",
], Inches(6.75), Inches(1.85), Inches(6.2), Inches(0.4), size=11, spacing=0.44, color=BLACK)

callout(s,
        "The ceiling is NOT a modelling failure — it is a DATA failure.\n"
        "Missing columns (GHSV, CH₄:O₂, pressure) explain why Y>15% catalysts exist.\n"
        "No algorithm can compensate for absent input features.",
        Inches(6.7), Inches(4.0), Inches(6.3), Inches(1.1),
        fill=RGBColor(0xE8, 0xF5, 0xFF), border=ACCENT,
        text_color=DGRAY, size=11, italic=False)

add_text(s,
         "Fix the data → the model automatically improves.  That is the single highest-impact next step.",
         Inches(6.7), Inches(5.25), Inches(6.3), Inches(0.55),
         size=12, bold=True, color=NAVY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Statistical Rigour & Honesty (Chapter 8/9)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Statistical Rigour & Honesty", chapter="Chapter 8")
nav_bar(s, 13, TOTAL)

# Left: the headline is coverage-inflated
add_text(s, "The headline was coverage-inflated",
         Inches(0.3), Inches(1.2), Inches(6.4), Inches(0.4),
         size=13, bold=True, color=NAVY)
bullets(s, [
    "Cells run further contain better yields, so a score over all",
    "917 catalysts is partly a record of which runs were finished.",
    "All 917:        Spearman 0.767   enrichment 4.35×",
    "Equal effort (771):  Spearman 0.724   enrichment 3.77×",
    "Coupling removed BY MEASUREMENT: ρ(n_rows, max) +0.293 → +0.003",
    "300 random 771-subsets give 0.767 [0.756, 0.780] — 0.724 sits",
    "below that interval, so the drop is real, not a smaller-sample effect.",
], Inches(0.35), Inches(1.7), Inches(6.3), Inches(0.4), size=11, spacing=0.42, color=BLACK)

# Right: the negative control that justifies enrichment as primary
add_text(s, "The sharpest test we could devise",
         Inches(6.9), Inches(1.2), Inches(6.1), Inches(0.4),
         size=13, bold=True, color=NAVY)
bullets(s, [
    "We retrained the model to predict HOW MANY measurements a",
    "catalyst received. It never sees a yield.",
    "That model reaches Spearman 0.400 against observed max yield…",
    "…but its enrichment is 0.87× — no better than random.",
    "Rank correlation is partly purchasable from experimental effort.",
    "Enrichment is not. That is why enrichment is the primary metric.",
], Inches(6.95), Inches(1.7), Inches(6.0), Inches(0.4), size=11, spacing=0.42, color=BLACK)

callout(s,
        "And the open question is bounded: ranking catalysts by their observed FLOOR instead of their "
        "ceiling shares only 7 of the top 20, but a ceiling-trained model loses just 0.017 Spearman "
        "against floor ground truth — less than seed-averaging alone buys — and never drops below "
        "4.02× enrichment. The answer changes the labels, not the decision.",
        Inches(0.35), Inches(5.15), Inches(12.6), Inches(1.5),
        fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN,
        text_color=GDARK, size=11)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Critical Analysis — Per-Method Review (Chapter 9a)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Critical Analysis — Per-Method Review", chapter="Chapter 9")
nav_bar(s, 14, TOTAL)

# Left: per-method critique table
add_text(s, "Honest critique of each method",
         Inches(0.4), Inches(1.35), Inches(7.0), Inches(0.4),
         size=14, bold=True, color=NAVY)
add_rect(s, Inches(0.4), Inches(1.82), Inches(7.0), Inches(0.4), fill=NAVY)
add_text(s, "Method", Inches(0.5), Inches(1.85), Inches(1.8), Inches(0.34),
         size=11, bold=True, color=WHITE)
add_text(s, "Main limitation", Inches(2.35), Inches(1.85), Inches(4.9), Inches(0.34),
         size=11, bold=True, color=WHITE)
crit_rows = [
    ("Naive merge", "Negative control — makes the model worse", False),
    ("DRST", "Hard cliff; discards ~80%; τ tuned on the reported CV", False),
    ("KMM", "O(n²) kernel; σ-sensitive; labels still in loss; ties DRST", False),
    ("Prior FT", "Best — Stage-1 on 782 rows → high variance; two models", True),
]
y = 2.22
for meth, lim, best in crit_rows:
    h = 0.58
    bg = RGBColor(0xE8, 0xF5, 0xE9) if best else LIGHT
    add_rect(s, Inches(0.4), Inches(y), Inches(7.0), Inches(h), fill=bg, line=LGRAY)
    add_text(s, meth, Inches(0.5), Inches(y + 0.05), Inches(1.8), Inches(h - 0.08),
             size=10, bold=True, color=GDARK if best else BLACK)
    add_text(s, lim, Inches(2.35), Inches(y + 0.05), Inches(4.9), Inches(h - 0.08),
             size=10, color=GDARK if best else DGRAY)
    y += h + 0.04

# Right: cross-cutting limits
add_rect(s, Inches(7.7), Inches(1.35), Inches(5.3), Inches(0.42), fill=ACCENT)
add_text(s, "Cross-cutting limits",
         Inches(7.85), Inches(1.38), Inches(5.0), Inches(0.36),
         size=13, bold=True, color=WHITE)
bullets(s, [
    "All methods evaluated on our-data CV only; OOD tested for Prior FT alone",
    "Point predictions — no uncertainty estimates to guide experiment selection",
    "Y>15% ceiling is DATA-bound (missing GHSV, CH₄:O₂, pressure), not model-bound",
], Inches(7.75), Inches(1.88), Inches(5.2), Inches(0.4), size=10, spacing=0.42, color=BLACK)

callout(s,
        "Key insight: the ceiling effect and much of the label shift share the same root cause — "
        "missing operating-condition features. Fix the data; the models follow.",
        Inches(7.7), Inches(3.35), Inches(5.3), Inches(0.95),
        fill=RGBColor(0xE8, 0xF5, 0xFF), border=ACCENT,
        text_color=DGRAY, size=10, italic=False)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — What's Next — Priority-Ordered (Chapter 9b)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "What's Next — Priority-Ordered", chapter="Chapter 9")
nav_bar(s, 15, TOTAL)

next_steps = [
    ("1", "Ask JAIST for the reaction-condition columns",
     "One email. Converts 19.9% of currently unreachable variance into modellable signal and turns "
     "917 training examples back into 89,074.", True),
    ("2", "Ask whether the 27 slots are conditions or time-on-stream",
     "Provably unanswerable from the file: every cell is stored sorted by yield, so row order records "
     "rank, not acquisition sequence. Only the lab can settle it.", True),
    ("3", "Ask why grid coverage is incomplete",
     "186 catalyst-temperature cells absent; only 811 of 917 catalysts have all five temperatures. "
     "Decides whether the bias is correctable or itself informative.", False),
    ("4", "Re-scope the campaign before reactor time is spent",
     "20 runs per catalyst instead of 135 reproduces the full ranking at ρ = 0.949, buying ~72 "
     "catalysts plus a control arm for the budget of 17 exhaustive ones.", False),
    ("5", "Send the corrected work note",
     "Complete and verified. Every number now traces to a JSON written by a committed script.", False),
    ("6", "Run the prospective validation",
     "Prof. Taniike offered to synthesise candidates. A shortlist is a set to test, not a league "
     "table — the model cannot order within its own top 20.", False),
]

col1 = next_steps[:3]
col2 = next_steps[3:]

for col_idx, col_items in enumerate([col1, col2]):
    x_left = Inches(0.4) if col_idx == 0 else Inches(6.85)
    col_w = Inches(6.2)
    y = 1.45
    for num, head, body, top in col_items:
        bg = RGBColor(0xE8, 0xF5, 0xE9) if top else LIGHT
        border_c = GREEN if top else LGRAY
        add_rect(s, x_left, Inches(y), col_w, Inches(1.55), fill=bg, line=border_c)
        badge(s, num, x_left + Inches(0.1), Inches(y + 0.1),
              Inches(0.36), Inches(0.36),
              fill=GREEN if top else ACCENT, size=14)
        add_text(s, head, x_left + Inches(0.55), Inches(y + 0.1),
                 col_w - Inches(0.65), Inches(0.42),
                 size=12, bold=True, color=GDARK if top else NAVY)
        add_text(s, body, x_left + Inches(0.55), Inches(y + 0.55),
                 col_w - Inches(0.65), Inches(0.85),
                 size=10, color=DGRAY)
        if top:
            badge(s, "Highest impact", x_left + col_w - Inches(1.85), Inches(y + 0.1),
                  Inches(1.7), Inches(0.3), GREEN, size=9)
        y += 1.65

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Summary
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_title_bar(s, "Summary")
nav_bar(s, 16, TOTAL)

# Big summary box
add_rect(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(4.6),
         fill=LIGHT, line=NAVY)
add_text(s,
         "What we asked",
         Inches(1.25), Inches(1.55), Inches(10.9), Inches(0.4),
         size=14, bold=True, color=NAVY)
add_text(s,
         "Can published OCM literature improve a lab-data model — and can we rank UNSEEN catalysts well enough to guide synthesis?",
         Inches(1.25), Inches(1.95), Inches(10.9), Inches(0.45),
         size=12, color=DGRAY)

add_text(s,
         "What we learned",
         Inches(1.25), Inches(2.6), Inches(10.9), Inches(0.4),
         size=14, bold=True, color=NAVY)
findings = [
    "The original −9.7% gain was catalyst-identity leakage. Under catalyst-grouped CV the same models are 1.8% WORSE than baseline.",
    "No literature-integration design beats composition alone in-domain — four pre-registered designs and a 28-family follow-up, all null.",
    "Literature helps only where the lab has NO coverage: on non-impregnation chemistry a plain merge lifts Spearman 0.24 → 0.39 and enrichment 0.42× → 1.34×. Plain merging beats the two-stage machinery — the value is the data, not the method.",
    "The composition-only model still screens unseen catalysts usefully: enrichment 3.77× on the equal-effort set (95% CI 3.04–4.89×), and that conclusion survives either reading of the open data question.",
]
for i, txt in enumerate(findings):
    add_text(s, "•  " + txt,
             Inches(1.35), Inches(3.05) + Inches(i * 0.4),
             Inches(10.7), Inches(0.4),
             size=12, color=BLACK)

# Bottom emphasis box
add_rect(s, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.85),
         fill=RGBColor(0xE8, 0xF5, 0xE9), line=GREEN)
add_text(s, "Next step that matters most",
         Inches(1.25), Inches(6.23), Inches(11.0), Inches(0.4),
         size=13, bold=True, color=GDARK)
add_text(s,
         "Ask for the reaction-condition columns. One email unlocks more than any modelling change we have tried.",
         Inches(1.25), Inches(6.6), Inches(11.0), Inches(0.4),
         size=12, color=GDARK)

# ── Save ──────────────────────────────────────────────────────────────────────
pptx_path = ROOT / "ocm_presentation.pptx"
prs.save(str(pptx_path))
print(f"Saved: {pptx_path}  ({pptx_path.stat().st_size // 1024} KB)")

# ── Export to PDF via LibreOffice ─────────────────────────────────────────────
import subprocess, shutil
lo = shutil.which("libreoffice") or shutil.which("soffice")
pdf = ROOT / "ocm_presentation.pdf"
if lo:
    print("Converting to PDF via LibreOffice...")
    # Delete first, then require a NEW file. The previous version reported success whenever the
    # path merely existed, so a failed conversion silently left the STALE pdf in place and still
    # printed "Saved" -- which is exactly how a retracted claim survived in ocm_presentation.pdf
    # after the deck was corrected. libreoffice-impress must be installed; libreoffice-core alone
    # cannot load .pptx and exits 0 while printing "source file could not be loaded".
    if pdf.exists():
        pdf.unlink()
    result = subprocess.run(
        [lo, "--headless", "--convert-to", "pdf",
         "--outdir", str(ROOT), str(pptx_path)],
        capture_output=True, text=True, timeout=300
    )
    combined = (result.stdout or "") + (result.stderr or "")
    if not pdf.exists() or "could not be loaded" in combined:
        raise SystemExit(
            "LibreOffice PDF export FAILED (no file produced).\n"
            "Install libreoffice-impress: libreoffice-core alone cannot read .pptx.\n"
            + combined[:800])
    print(f"Saved: {pdf}  ({pdf.stat().st_size // 1024} KB)")
else:
    raise SystemExit("libreoffice/soffice not found - cannot export PDF")
